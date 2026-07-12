"""Build the Strata Story PowerPoint deck (16:9, navy + gold brand)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GOLD = RGBColor(0xA6, 0x8A, 0x64)
AMBER = RGBColor(0xC0, 0x70, 0x40)
LIGHT = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)
BG_LIGHT = RGBColor(0xFA, 0xF7, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multi_text(slide, left, top, width, height, items, anchor=MSO_ANCHOR.TOP):
    """items = list of dicts: {text, size, bold, color, align, italic}."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = it.get('align', PP_ALIGN.LEFT)
        if 'space_after' in it:
            p.space_after = Pt(it['space_after'])
        r = p.add_run()
        r.text = it['text']
        r.font.name = 'Calibri'
        r.font.size = Pt(it.get('size', 16))
        r.font.bold = it.get('bold', False)
        r.font.italic = it.get('italic', False)
        r.font.color.rgb = it.get('color', NAVY)
    return tb


def add_footer(slide, page_num, total):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), NAVY)
    add_text(slide, Inches(0.3), SH - Inches(0.32), Inches(8), Inches(0.3),
             'THE STRATA STORY · 369 Alliance', size=9, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.32), Inches(0.9), Inches(0.3),
             f'{page_num} / {total}', size=9, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header_bar(slide, chapter_num, chapter_title):
    add_rect(slide, 0, 0, SW, Inches(0.85), NAVY)
    add_rect(slide, 0, Inches(0.85), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.4), Inches(0.1), Inches(2.2), Inches(0.7),
             f'CHAPTER {chapter_num}', size=12, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2.6), Inches(0.1), Inches(10.5), Inches(0.7),
             chapter_title, size=24, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, left, top, width, height, data, header_fill=NAVY, header_color=WHITE,
              row_fills=None, col_widths=None, row_height=None, font_size=12, header_size=12):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w
    if row_height:
        for ri in range(rows):
            tbl.rows[ri].height = row_height
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = 'Calibri'
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                r.font.bold = True
                r.font.size = Pt(header_size)
                r.font.color.rgb = header_color
            else:
                cell.fill.solid()
                if row_fills:
                    cell.fill.fore_color.rgb = row_fills[ri % len(row_fills)]
                else:
                    cell.fill.fore_color.rgb = BG_LIGHT if ri % 2 == 0 else WHITE
                r.font.size = Pt(font_size)
                r.font.color.rgb = NAVY
    return tbl_shape


def add_rule_box(slide, top, rule, who, law, height=Inches(1.6)):
    left = Inches(0.6)
    width = SW - Inches(1.2)
    add_rect(slide, left, top, width, height, LIGHT)
    add_rect(slide, left, top, Inches(0.15), height, GOLD)
    items = [
        {'text': 'Rule in one line: ', 'size': 13, 'bold': True, 'color': NAVY},
        {'text': rule, 'size': 13, 'color': NAVY},
    ]
    # Multi-paragraph
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.1), width - Inches(0.4), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = 'Rule in one line: '
    r.font.name = 'Calibri'; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = rule
    r2.font.name = 'Calibri'; r2.font.size = Pt(13); r2.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r = p2.add_run(); r.text = 'Who pays: '
    r.font.name = 'Calibri'; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = AMBER
    r2 = p2.add_run(); r2.text = who
    r2.font.name = 'Calibri'; r2.font.size = Pt(12); r2.font.color.rgb = NAVY

    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)
    r = p3.add_run(); r.text = 'The law: '
    r.font.name = 'Calibri'; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GOLD
    r2 = p3.add_run(); r2.text = law
    r2.font.name = 'Calibri'; r2.font.size = Pt(12); r2.font.italic = True; r2.font.color.rgb = NAVY


def add_bullets(slide, left, top, width, height, items, size=14, color=NAVY, bullet='•'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f'{bullet}  {it}'
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.color.rgb = color


# ====================================================================
# Slide builders
# ====================================================================

slides_to_build = []  # we'll call functions, then renumber footers


def s_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.0), SW, Inches(0.08), GOLD)
    add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.2),
             'THE STRATA STORY', size=64, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.7),
             "A simple guide for owners — what's mine, what's \"ours\", and who pays when something breaks",
             size=20, italic=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
             'Based on NSW Strata Legislation', size=14, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.85), Inches(11.5), Inches(0.5),
             'Building Commission NSW · 369 Alliance', size=12, italic=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
             '13 chapters · 1 cheat sheet · plain English', size=11, color=WHITE)
    return s


def s_intro():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '0', 'How to read this guide')
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             'This is a story, not a law book.', size=22, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.6),
             'Every chapter answers ONE question owners ask all the time.', size=16, color=GREY)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.4),
             'At the end of every chapter you will find:', size=15, bold=True, color=AMBER)
    add_bullets(s, Inches(0.9), Inches(3.15), Inches(12), Inches(2),
                ['Rule in one line.', 'Who pays.', 'The law behind it.'], size=15)
    add_rect(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5), LIGHT)
    add_rect(s, Inches(0.6), Inches(5.0), Inches(0.15), Inches(1.5), GOLD)
    add_text(s, Inches(1.0), Inches(5.15), Inches(11.5), Inches(0.5),
             'If you only have 5 minutes…', size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.65), Inches(11.5), Inches(0.8),
             'Read CHAPTER 2 (the imaginary line) and CHAPTER 5 (the decision tree). '
             'Those two answer 80% of the calls we get.', size=14, color=NAVY)
    return s


def s_ch1():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '1', 'What is Strata, really?')
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.7),
             'Imagine 50 families in one building.', size=24, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(1.0),
             'Each family wants their own front door, kitchen, bedroom — but everyone has to share the lift, the roof, the driveway, the walls between apartments, and the pipes inside the walls.',
             size=15, color=GREY)
    add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.5),
             'And everyone shares the cost of fixing those shared things — forever.', size=15, italic=True, color=AMBER)
    add_text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.6),
             'That is Strata.', size=22, bold=True, color=NAVY)

    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
             'When you buy an apartment in NSW, you buy TWO things at once:', size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(4.85), Inches(12), Inches(1.0),
                ['Your Lot — the box of air inside your apartment.',
                 'A share of the Common Property — a slice of everything that is shared.'],
                size=14)
    add_rule_box(s, Inches(5.85),
                 'Buying an apartment means buying a private box + a share of a shared building.',
                 'You pay levies (your share of the bills) every quarter, for as long as you own.',
                 'Strata Schemes Management Act 2015 (SSMA) and Strata Schemes Development Act 2015 (SSDA).',
                 height=Inches(1.25))
    return s


def s_ch2_a():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '2', 'The Imaginary Line')
    add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.7),
             'There is a line that runs around the inside of every apartment.', size=20, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
             "You can't see it — but it's drawn on the registered Strata Plan.", size=14, italic=True, color=GREY)

    # Two columns
    # Left: LOT
    add_rect(s, Inches(0.6), Inches(2.7), Inches(6.0), Inches(0.5), AMBER)
    add_text(s, Inches(0.7), Inches(2.74), Inches(5.9), Inches(0.45),
             'YOUR side of the line  =  LOT PROPERTY  =  YOU', size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(3.5), BG_LIGHT)
    add_bullets(s, Inches(0.8), Inches(3.3), Inches(5.7), Inches(3.4),
                ['Light fittings, switches, power points',
                 'Blinds, curtains, internal doors',
                 'Internal walls NOT shown on the Strata Plan',
                 'Floor finishes (carpet, your tiles)',
                 'Kitchen cupboards & appliances',
                 'Tiles & waterproofing on internal walls',
                 'Bath, basin, toilet, taps',
                 'Plumbing ABOVE the upper surface of the floor'],
                size=12)

    # Right: COMMON
    add_rect(s, Inches(6.85), Inches(2.7), Inches(6.0), Inches(0.5), NAVY)
    add_text(s, Inches(6.95), Inches(2.74), Inches(5.9), Inches(0.45),
             'OTHER side  =  COMMON PROPERTY  =  OWNERS CORP', size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(6.85), Inches(3.2), Inches(6.0), Inches(3.5), BG_LIGHT)
    add_bullets(s, Inches(7.05), Inches(3.3), Inches(5.7), Inches(3.4),
                ['The structure (slabs, ceilings, columns, beams)',
                 'External windows, doors, balcony doors',
                 'ORIGINAL floor tiles (incl. shower)',
                 'ORIGINAL waterproofing membrane',
                 'Tiles & waterproofing on BOUNDARY walls',
                 'Shared plumbing inside walls/voids',
                 'Wiring outside your apartment / shared',
                 'Front door of your apartment (on the plan)'],
                size=12)
    return s


def s_ch2_b():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '2', 'The Imaginary Line — applies to all schemes since 1 July 1974')
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             'The 3 questions that decide everything:', size=20, bold=True, color=NAVY)
    # Three boxes
    questions = [
        ('1', 'Where is it?', 'Inside or outside the line on the Strata Plan?'),
        ('2', 'Was it there from day one?', 'Original = usually Common Property.'),
        ('3', 'Does it serve more than one apartment?', 'Yes = Common Property.'),
    ]
    for i, (n, q, a) in enumerate(questions):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(2.1), Inches(3.95), Inches(2.6), LIGHT)
        add_rect(s, x, Inches(2.1), Inches(3.95), Inches(0.55), NAVY)
        add_text(s, x + Inches(0.15), Inches(2.13), Inches(0.5), Inches(0.5),
                 n, size=22, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), Inches(2.13), Inches(3.0), Inches(0.5),
                 q, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(2.85), Inches(3.65), Inches(1.7),
                 a, size=13, color=NAVY)

    add_rule_box(s, Inches(5.0),
                 'If it was built originally and serves more than one owner, it is Common Property.',
                 'Lot side = the owner. Common side = all owners (through levies).',
                 'Strata Schemes Development Act 2015, Schedule 2 + your Strata Plan.',
                 height=Inches(2.0))
    return s


def s_ch3():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '3', 'The Cast — who is who in your building')
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             'Knowing who does what saves you years of frustration.', size=14, italic=True, color=GREY)
    data = [
        ['Character', 'What they actually do', 'Can they sign remedial works?'],
        ['Owner (you)', 'Owns 1 Lot + share of Common Property. Pays levies. Votes.', 'Only inside your own Lot.'],
        ['Owners Corporation (OC)', 'Legal "company" of ALL owners. Owns the Common Property.', 'YES — only body that can.'],
        ['Strata Committee', '3–9 owners elected for day-to-day decisions.', 'Limited — big jobs go to the OC.'],
        ['Strata Manager', 'Licensed agent for admin, levies, meetings, insurance.', 'NO — needs OC approval.'],
        ['Building Manager', 'Looks after the building (cleaning, gardens, small fixes).', 'NO — needs OC approval.'],
        ['NCAT (Tribunal)', 'Resolves disputes the OC cannot fix.', 'Decides; orders works in some cases.'],
        ['NSW Fair Trading', 'Free mediation before NCAT.', 'No — they mediate.'],
        ['Building Commission NSW', 'Regulator. Inspects works. Issues fines.', 'They police; not your project.'],
    ]
    add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.4), data,
              col_widths=[Inches(2.6), Inches(5.5), Inches(4.0)], font_size=11, header_size=12)
    add_rule_box(s, Inches(6.15),
                 'Only the Owners Corporation can authorise major remedial works. Not the strata manager. Not the building manager. Not one loud owner.',
                 'OC pays from the Capital Works Fund (or by special levy).',
                 'SSMA 2015 sections 5, 8, 9 and 106.',
                 height=Inches(1.0))
    return s


def s_ch4():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '4', 'The Rule Book — 8 laws that govern your apartment')
    data = [
        ['When this happens…', 'The law you point at'],
        ['Meetings, levies, by-laws, voting, common property repairs', 'Strata Schemes Management Act 2015 + Reg 2016'],
        ['The strata plan, subdivision, ending the scheme', 'Strata Schemes Development Act 2015 + Reg 2016'],
        ['Designs/builders for waterproofing, cladding, fire, structure (Class 2)', 'Design and Building Practitioners Act 2020 + Reg 2021'],
        ['Defects in residential apartment buildings', 'Residential Apartment Buildings Act 2020 (RAB)'],
        ['Builder warranties (2-year minor / 6-year major)', 'Home Building Act 1989'],
        ['Whether the work needs council approval (DA, CDC, exempt)', 'EP&A Act 1979 + Codes SEPP 2008'],
        ["Strata manager's licence and conduct", 'Property and Stock Agents Act 2002'],
        ['Buying or selling the apartment', 'Conveyancing Act 1919'],
    ]
    add_table(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(4.7), data,
              col_widths=[Inches(6.5), Inches(5.6)], font_size=12)
    add_rule_box(s, Inches(6.05),
                 'Day-to-day = SSMA. Building quality = DBP + RAB. Defects warranty = HBA. Approval needed? = EP&A + Codes SEPP.',
                 'Depends on the matter — see other chapters.',
                 'See each Act above.',
                 height=Inches(1.1))
    return s


def s_ch5():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '5', 'The Decision Tree — whose problem is it?')
    add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
             'Use this BEFORE you call anyone.', size=18, bold=True, color=NAVY)

    # Decision flow with rectangles + arrows (simplified visual)
    # Top: question box
    add_rect(s, Inches(4.0), Inches(2.0), Inches(5.3), Inches(0.8), NAVY)
    add_text(s, Inches(4.0), Inches(2.0), Inches(5.3), Inches(0.8),
             'Is the broken thing on the LOT side of the line?', size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Yes branch (left)
    add_text(s, Inches(2.5), Inches(3.05), Inches(1.5), Inches(0.4),
             'YES →', size=13, bold=True, color=AMBER, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.6), Inches(3.5), Inches(3.5), Inches(1.0), AMBER)
    add_text(s, Inches(0.6), Inches(3.5), Inches(3.5), Inches(1.0),
             'YOU pay.\nYOU fix.', size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # No branch (right) → second question
    add_text(s, Inches(9.4), Inches(3.05), Inches(1.5), Inches(0.4),
             '← NO', size=13, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(9.2), Inches(3.5), Inches(3.6), Inches(0.8), NAVY)
    add_text(s, Inches(9.2), Inches(3.5), Inches(3.6), Inches(0.8),
             'Is it ORIGINAL to the building?', size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # YES original
    add_rect(s, Inches(9.2), Inches(4.5), Inches(3.6), Inches(0.85), GOLD)
    add_text(s, Inches(9.2), Inches(4.5), Inches(3.6), Inches(0.85),
             'OWNERS CORP\npays & fixes', size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # NO not original
    add_rect(s, Inches(9.2), Inches(5.4), Inches(3.6), Inches(0.85), GREY)
    add_text(s, Inches(9.2), Inches(5.4), Inches(3.6), Inches(0.85),
             'Whoever installed it pays', size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(5.1), Inches(8.0), Inches(1.5),
             'Lot side examples: light fittings, your dishwasher, tiles you renovated yourself, plumbing above the floor.\n\nCommon side examples: original waterproofing, structure, external windows, shared pipes.',
             size=12, color=NAVY)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             'Rule:  Original + serves more than one Lot = Owners Corporation.   Yours alone, installed by you, above the floor = you.',
             size=12, bold=True, italic=True, color=AMBER)
    return s


def s_ch6_a():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '6', 'The Remediation Story — when something BIG breaks')
    add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.6),
             '"Big" = roof leaks across 6 apartments · façade cracking · basement waterproofing failed · cladding has to come off.',
             size=13, italic=True, color=GREY)
    add_text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5),
             'These are remedial building works — and they have their own rulebook now (since 2020).',
             size=14, bold=True, color=NAVY)

    steps = [
        ('1', 'Spot it', 'Owner / building manager / inspector reports it.'),
        ('2', 'Get advice', 'Planning Consultant says: Exempt / CDC / DA.'),
        ('3', 'OC VOTES', 'Only the Owners Corporation can authorise.'),
        ('4', 'DBP twist', 'Class 2 → Reg. Designer + Reg. Builder may be needed.'),
        ('5', 'Lodge designs', 'Upload to NSW Planning Portal BEFORE site work.'),
        ('6', 'Build', 'Build to the lodged designs. Variations? Re-lodge.'),
        ('7', 'Close out', 'Builder lodges Building Compliance Declaration.'),
    ]
    for i, (n, t, d) in enumerate(steps):
        col = i % 4
        row = i // 4
        x = Inches(0.6 + col * 3.15)
        y = Inches(2.55 + row * 2.05)
        add_rect(s, x, y, Inches(3.0), Inches(1.85), LIGHT)
        add_rect(s, x, y, Inches(3.0), Inches(0.45), NAVY)
        add_text(s, x + Inches(0.1), y + Inches(0.02), Inches(0.4), Inches(0.4),
                 n, size=18, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.55), y + Inches(0.02), Inches(2.4), Inches(0.4),
                 t, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.7), Inches(1.25),
                 d, size=11, color=NAVY)
    return s


def s_ch6_b():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '6', 'Remediation — the two DBP exceptions that catch everyone out')
    data = [
        ['Type of work', 'DBP triggered?', 'When can you avoid it?'],
        ['Waterproofing', 'YES (almost always)', 'Only if ALL: (1) bathroom/kitchen/laundry/toilet alteration, (2) one single dwelling, (3) exempt development.'],
        ['External cladding', 'YES (almost always)', 'Only if NOT being done because of a council order or development control order.'],
    ]
    add_table(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(2.6), data,
              col_widths=[Inches(2.6), Inches(3.0), Inches(6.5)], font_size=12)

    add_rect(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.0), AMBER)
    add_text(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(1.0),
             'Even if your work is "exempt development" under the Codes SEPP — waterproofing and cladding STILL need a Registered Design Practitioner.',
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_rule_box(s, Inches(5.4),
                 'Owners vote → designer designs → portal upload → builder builds → builder declares → case closed.',
                 'OC pays (capital works fund or special levy).',
                 'DBP Act 2020, DBP Regulation 2021 (clause 13), RAB Act 2020, Codes SEPP 2008.',
                 height=Inches(1.7))
    return s


def s_ch7():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '7', 'The "It\'s an Emergency" Exception')
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             'Sometimes you cannot wait for designs. The roof is open. People can\'t sleep. Real risk to safety.',
             size=14, italic=True, color=GREY)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
             'That is Emergency Remedial Building Work. The rules bend — but only a little.',
             size=15, bold=True, color=NAVY)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             'ALL 5 of these must be true (a "reasonable excuse"):', size=14, bold=True, color=AMBER)
    add_bullets(s, Inches(0.9), Inches(3.05), Inches(12), Inches(2.5),
                ['Immediate action is needed.',
                 'The issue is causing or likely to cause damage.',
                 'It impacts: ability to live there, OR health and safety, OR risk of further damage.',
                 'The impact is serious.',
                 'The work done is only what is necessary to mitigate it — not the full final fix.'],
                size=13)
    add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.7), GOLD)
    add_text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.7),
             'Then within 7 DAYS the builder lodges the emergency work + Building Compliance Declaration on the NSW Planning Portal.',
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rule_box(s, Inches(6.25),
                 'Emergency = act now, declare within 7 days, then plan the proper fix.',
                 'OC pays (capital works fund or special levy).',
                 'DBP Act 2020 — emergency remedial building work provisions.',
                 height=Inches(0.9))
    return s


def s_ch8():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '8', 'The Penalty Story — why this is not optional')
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             '"Can we just skip the paperwork?"  No. Here is what happens.', size=15, italic=True, color=GREY)
    data = [
        ['What went wrong', 'Body corporate', 'Individual'],
        ['Designer failed to provide design compliance declarations', 'Up to $165,000', 'Up to $55,000'],
        ['Builder failed to lodge declared designs', 'Up to $33,000', 'Up to $11,000'],
        ['Significant offences', 'Prosecution', 'Prosecution'],
    ]
    add_table(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.0), data,
              col_widths=[Inches(6.6), Inches(2.7), Inches(2.8)], font_size=13, header_size=13)

    add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
             'On top of fines, Building Commission NSW can issue:', size=13, bold=True, color=AMBER)
    add_bullets(s, Inches(0.9), Inches(4.65), Inches(12), Inches(1.5),
                ['Stop Work Orders (RAB Act).',
                 'Building Work Rectification Orders — forcing the OC to undo and redo.',
                 'Prohibition Orders — blocking Occupation Certificate or strata plan registration.'],
                size=13)
    add_rule_box(s, Inches(6.2),
                 'Cutting corners on DBP costs more than doing it properly.',
                 'OC and/or practitioners.',
                 'DBP Act 2020 + RAB Act 2020.',
                 height=Inches(0.9))
    return s


def s_ch9():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '9', 'But what about MY renovation?')
    add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
             'SSMA 2015 sorts your work into 3 buckets (sections 108–110):', size=15, italic=True, color=GREY)
    data = [
        ['Bucket', 'What it means', 'What you need'],
        ['Cosmetic Work (s.109)', 'Painting, hanging pictures, curtains, replacing handles, hooks, TV.', 'Nothing. Just go ahead.'],
        ['Minor Renovations (s.110)', 'New kitchen, built-in wardrobe, internal walls, recessed lights, hardwood/laminate flooring.', 'OC approval by ordinary resolution (50%+).'],
        ['Major Renovations', 'Anything that changes external appearance, OR touches structure / waterproofing / fire / common property.', 'By-law by special resolution (75%+) — usually with DBP-compliant designer + builder.'],
    ]
    add_table(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(3.0), data,
              col_widths=[Inches(2.4), Inches(5.7), Inches(4.0)], font_size=12)
    add_rect(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.0), AMBER)
    add_text(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(1.0),
             'TRAP: Replacing your bathroom waterproofing membrane is NOT cosmetic. It touches the original waterproofing — which is Common Property. OC approval + Registered Design Practitioner needed.',
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rule_box(s, Inches(6.15),
                 'Changes the look from outside, structure, fire, or waterproofing? It is a Major Renovation.',
                 'You pay for your own work; by-law often makes you pay ongoing maintenance too.',
                 'SSMA 2015 sections 108–110, plus DBP Act 2020 for waterproofing.',
                 height=Inches(0.95))
    return s


def s_ch10():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '10', 'The First-Year Story — a brand-new building')
    data = [
        ['Moment', 'What happens'],
        ['Strata Plan registered', 'Owners Corporation is born. Developer holds all votes (owns unsold lots).'],
        ['Initial Period', 'Until 1/3 of lots are sold to outside buyers. Developer cannot make big Common Property decisions.'],
        ['First AGM', 'Within 2 months of end of Initial Period. Owners elect first Strata Committee. Developer hands over plans, manuals, maintenance schedule.'],
        ['Building Bond (Class 2, 4+ storeys)', 'Developer lodges 2% of contract price with NSW Fair Trading.'],
        ['Inspections', 'Independent Building Inspector: interim (15–18 months) + final (21–24 months).'],
        ['HBA Defects Warranty', '2-year minor / 6-year major from completion (Home Building Act s.18B).'],
        ['DBP Statutory Duty of Care', '10-year duty of care from designer, builder, manufacturer, supplier (DBP Act s.37). Retrospective to 11 June 2010.'],
    ]
    add_table(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(4.5), data,
              col_widths=[Inches(3.5), Inches(8.6)], font_size=11, header_size=12)
    add_rule_box(s, Inches(5.85),
                 'Year 1–2 = developer hand-over and bond. Year 0–10 = DBP duty of care. Year 0–6 = HBA major defects.',
                 'Developer/builder for warranty work; OC for ongoing maintenance after handover.',
                 'SSMA 2015 Pt 3, SSDA 2015, DBP Act s.37, HBA s.18B + Pt 11.',
                 height=Inches(1.25))
    return s


def s_ch11():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '11', 'How to raise an issue — the ladder')
    rungs = [
        ('1', 'Speak to your Strata Committee Secretary or your Strata Manager (in writing).'),
        ('2', 'Strata Committee tries to resolve within delegated authority.'),
        ('3', 'If unresolved → raise as a Motion at the next OC general meeting.'),
        ('4', 'If unresolved → FREE mediation with NSW Fair Trading.'),
        ('5', 'Final step → Apply to NCAT (Strata & Community Schemes Division).'),
    ]
    for i, (n, t) in enumerate(rungs):
        y = Inches(1.4 + i * 0.85)
        add_rect(s, Inches(0.6), y, Inches(0.7), Inches(0.7), NAVY)
        add_text(s, Inches(0.6), y, Inches(0.7), Inches(0.7),
                 n, size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.3), y, Inches(11.4), Inches(0.7), LIGHT)
        add_text(s, Inches(1.5), y, Inches(11.2), Inches(0.7),
                 t, size=13, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_rule_box(s, Inches(6.2),
                 'Always escalate one step at a time. Mediation is free. NCAT is the final step.',
                 'Each step has its own filing fee at NCAT.',
                 'SSMA 2015 Part 12, NCAT Civil and Administrative Tribunal Act 2013.',
                 height=Inches(1.0))
    return s


def s_ch12():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '12', 'The 7 Mistakes Owners Make')
    data = [
        ['#', 'The mistake', 'The fix'],
        ['1', '"It\'s my apartment, I can do what I want."', 'Read your Strata Plan + by-laws first.'],
        ['2', '"The Strata Manager said it\'s fine."', 'Strata Manager has NO authority to approve remedial works. Only the OC does.'],
        ['3', '"It\'s just waterproofing — no designer needed."', 'On common property = DBP. Skip designer = up to $165,000 fine.'],
        ['4', '"No planning consultant — it\'s just a balcony."', 'Planter removal, balustrade swap, balcony enlargement → DA + DBP.'],
        ['5', '"Pay cash, skip the Portal."', 'No Building Compliance Declaration = no closure. Future buyers walk.'],
        ['6', '"Defects are old, no point chasing."', 'DBP duty of care = 10 years. HBA major defects = 6 years.'],
        ['7', '"We voted last year — that\'s enough."', 'OC resolutions older than 12 months may need a re-vote.'],
    ]
    add_table(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.7), data,
              col_widths=[Inches(0.6), Inches(5.5), Inches(6.0)], font_size=11, header_size=12)
    return s


def s_ch13():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '13', 'The 30-Second Glossary')
    glossary = [
        ('OC', 'Owners Corporation'),
        ('SC', 'Strata Committee (3–9 elected owners)'),
        ('SM / BM', 'Strata Manager / Building Manager'),
        ('Common Property', 'Outside the Lot line, owned together'),
        ('Lot Property', 'Inside the Lot line, owned by you'),
        ('By-law', 'Rule that binds every owner, registered against SP'),
        ('Strata Plan (SP)', 'The legal map of the building'),
        ('BMS / SMS', 'Building / Strata Management Statement'),
        ('BMC', 'Building Management Committee'),
        ('DA / CDC / Exempt', 'The 3 planning approval pathways'),
        ('DBP / RAB', 'The two 2020 reform Acts'),
        ('CIRD', 'Construction-Issued Regulated Design'),
        ('DCD', 'Design Compliance Declaration'),
        ('BCD', 'Building Compliance Declaration'),
        ('PC / OC', 'Practical Completion / Occupation Certificate'),
        ('DLP', 'Defects Liability Period (usually 12 months)'),
        ('NCAT', 'NSW Civil and Administrative Tribunal'),
        ('NCC / BCA', 'National Construction Code / Building Code of Australia'),
    ]
    # 2 columns of glossary
    half = (len(glossary) + 1) // 2
    col_a = glossary[:half]
    col_b = glossary[half:]
    add_text(s, Inches(0.6), Inches(1.15), Inches(6), Inches(0.4),
             'Acronyms', size=14, bold=True, color=AMBER)
    add_text(s, Inches(6.85), Inches(1.15), Inches(6), Inches(0.4),
             'More acronyms', size=14, bold=True, color=AMBER)
    for i, (term, defn) in enumerate(col_a):
        y = Inches(1.6 + i * 0.5)
        add_text(s, Inches(0.6), y, Inches(2.2), Inches(0.45),
                 term, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.5), y, Inches(4.0), Inches(0.45),
                 defn, size=11, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    for i, (term, defn) in enumerate(col_b):
        y = Inches(1.6 + i * 0.5)
        add_text(s, Inches(6.85), y, Inches(2.2), Inches(0.45),
                 term, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.75), y, Inches(4.0), Inches(0.45),
                 defn, size=11, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    return s


def s_cheat_sheet():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, '★', 'The 1-Page Cheat Sheet — print this')
    add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
             'Whose problem is it?', size=15, bold=True, color=AMBER)
    data = [
        ['The thing', 'Lot or Common?', 'Who pays?'],
        ['Light bulb in your kitchen', 'LOT', 'YOU'],
        ['Dishwasher leak', 'LOT', 'YOU'],
        ['Original bathroom tiles cracked', 'COMMON', 'OC'],
        ['New tiles you laid leak', 'LOT', 'YOU'],
        ['Original waterproofing failed', 'COMMON', 'OC'],
        ['External window broken', 'COMMON', 'OC'],
        ['Roof leaking', 'COMMON', 'OC'],
        ['Pipe in wall, only your unit', 'LOT', 'YOU'],
        ['Pipe in wall, multiple units', 'COMMON', 'OC'],
        ['Front door of your apartment', 'COMMON', 'OC'],
        ['Lock on your front door', 'LOT', 'YOU'],
        ['Carpet inside your apartment', 'LOT', 'YOU'],
        ['Carpet in the corridor', 'COMMON', 'OC'],
    ]
    add_table(s, Inches(0.6), Inches(1.55), Inches(7.5), Inches(5.4), data,
              col_widths=[Inches(4.0), Inches(2.0), Inches(1.5)], font_size=10.5, header_size=11)

    # Right column: 4 questions
    add_rect(s, Inches(8.4), Inches(1.55), Inches(4.5), Inches(5.4), NAVY)
    add_text(s, Inches(8.6), Inches(1.7), Inches(4.2), Inches(0.5),
             'Before you start ANY work, ask:', size=14, bold=True, color=GOLD)
    questions = [
        '1.  Does it touch Common Property?\n     → Need OC approval.',
        '2.  Changes look, structure, fire or waterproofing?\n     → Need a by-law (75%+).',
        '3.  Waterproofing or cladding on Class 2?\n     → Need Reg. Designer + Reg. Builder.',
        '4.  Need a DA or CDC?\n     → Engage a Planning Consultant.',
    ]
    for i, q in enumerate(questions):
        add_text(s, Inches(8.6), Inches(2.3 + i * 1.0), Inches(4.2), Inches(0.95),
                 q, size=11, color=WHITE)
    add_text(s, Inches(8.6), Inches(6.4), Inches(4.2), Inches(0.45),
             'YES to ANY of 1–3 → it\'s NOT a DIY job.', size=12, bold=True, color=GOLD)
    return s


def s_closing():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.4), SW, Inches(0.08), GOLD)
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
             'Confused about Strata?', size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.7),
             'You are not alone — and you do not have to figure it out alone.', size=20, italic=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
             '369 Alliance · Building Compliance & Inspection — NSW', size=14, color=WHITE)
    add_text(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5),
             'Sources: SSMA 2015 · SSDA 2015 · DBP Act 2020 + Reg 2021 · RAB Act 2020 · HBA 1989 · Codes SEPP 2008 · EP&A Act 1979 · Building Commission NSW.',
             size=10, italic=True, color=GOLD)
    return s


# Build all
builders = [
    s_cover, s_intro,
    s_ch1, s_ch2_a, s_ch2_b, s_ch3, s_ch4, s_ch5,
    s_ch6_a, s_ch6_b, s_ch7, s_ch8, s_ch9, s_ch10,
    s_ch11, s_ch12, s_ch13, s_cheat_sheet, s_closing
]
slides = [b() for b in builders]
total = len(slides)

# Now add page footer to all slides except cover (slide 1) and closing
for idx, slide in enumerate(slides, start=1):
    if idx == 1 or idx == total:
        continue
    add_footer(slide, idx, total)

out = "C:/2026/369 Alliance/Website/docs/strata/storytelling/Strata_Story_for_Owners.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", total)
