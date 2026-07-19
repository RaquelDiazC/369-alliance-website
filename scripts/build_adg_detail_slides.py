"""
Add a 'detailed evidence breakdown' second page after each case slide
in ADG_6_Worked_Site_Exercise_Realistic_Drawings.pptx.

Visual style is matched to the Manus pages:
- Dark navy top strip + cream body + gold rule
- Right-rail evidence-board language transferred to a 4-panel Parts 1-4 grid
- Teal 'Design consequence' + dark red 'Performance Solution risk' style
  re-used at bottom for 'Design intent' + 'Approval evidence package'
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# --- Brand palette sampled from the existing slides ----------------------
NAVY      = RGBColor(0x14, 0x25, 0x36)   # top header strip
NAVY_DARK = RGBColor(0x20, 0x32, 0x44)   # footer
CREAM     = RGBColor(0xfe, 0xfc, 0xf8)   # body background
PAPER     = RGBColor(0xff, 0xff, 0xff)   # panel background
RULE      = RGBColor(0xd8, 0xcd, 0xb6)   # rule line
GOLD      = RGBColor(0xc9, 0xa9, 0x6d)   # pill border / accent
TEAL      = RGBColor(0x3e, 0x79, 0x7d)   # Design intent
RED       = RGBColor(0x9c, 0x5a, 0x56)   # Approval evidence
INK       = RGBColor(0x14, 0x25, 0x36)   # body text
INK_SOFT  = RGBColor(0x4a, 0x55, 0x60)   # subtle text
HEAD_GOLD = RGBColor(0xe8, 0xd8, 0xae)   # header text gold
WHITE     = RGBColor(0xff, 0xff, 0xff)


# --- Case data ----------------------------------------------------------
CASES = [
    {
        "n": "01",
        "title": "Inner West Flat Infill — Marrickville",
        "p1": [
            ("1A Type", "3-storey walk-up infill"),
            ("1B Character",
             "Federation cottages + 1970s walk-ups · brick + tile · "
             "intermittent canopy · 3–6 m gardens · low fences"),
            ("1C Precinct",
             "Individual-site response · character review on garden depth, "
             "fence type, brick palette"),
        ],
        "p2": [
            ("2A primary",  "3-storey + screened roof terrace inside 12 m datum"),
            ("2B envelope", "~6,200 m³ container · ~1,680 m² achievable GFA"),
            ("2C height",   "LEP 11 m / 3 storeys · effective height < 12 m"),
            ("2D FSR",      "1.1:1 ≈ 1,694 m² · courtyard option chosen"),
            ("2E depth",    "11.4 m glass-to-glass · dual-aspect corners"),
            ("2F separation", "12 / 9 / 6 m at < 12 m band"),
            ("2G street",   "4.5 m garden setback · low brick fence · retain street tree"),
            ("2H side/rear","3 m sides · 6 m rear lane · deep-soil reserve"),
            ("NCC overlay", "Type C · effective height < 12 m"),
        ],
        "p3": [
            "3A  tree register (2 jacarandas retained) + neighbour window survey",
            "3B  long axis E–W · north exposure · courtyard sun 9–15 mid-winter",
            "3C  front-loaded entry on Premier St + planted privacy buffer",
            "3D  rear north-facing courtyard ~280 m² (≥ 25 % site)",
            "3E  ~9 % deep soil · basement pulled 3 m off rear boundary",
            "3F  angled bedroom windows on south façade to #20",
            "3G  1:20 accessible route boundary to lift",
            "3H  driveway off rear lane only — no street crossover",
            "3J  0.5 sp/unit + 1 bike sp/unit · EV conduit only",
        ],
        "p4": [
            "4A  ≥ 70 % units 2 hr mid-winter sun · neighbour shadow study",
            "4B  65 % cross-vent via courtyard + corner units",
            "4C  2.7 m habitable · 2.4 m wet · services coordinated",
            "4D  furnished 1B / 2B / 2B+S / 3B · storage not from bedrooms",
            "4E  8 / 10 / 12 m² balconies · 2 m min depth · recessed at lane",
            "4G  tagged storage · 50 % inside / 50 % basement cages",
            "4H  wet stacks aligned · bedrooms off lift & chute",
            "4O  retained jacarandas + new canopy + driveway rain-garden",
            "4U  BASIX passive · roof PV",
        ],
        "design": (
            "Courtyard massing is preferred because it protects the street "
            "tree and brick-garden character while giving north-facing "
            "communal open space, solar access and genuine cross-ventilation "
            "to a majority of units."
        ),
        "approval": (
            "Site analysis · shadow study (mid-winter 9/12/15) · unit "
            "schedule with solar + ventilation tags · deep-soil plan · "
            "tree retention strategy · BASIX certificate · NCC Type C "
            "confirmation memo. No Performance Solution expected."
        ),
    },
    {
        "n": "02",
        "title": "Sloping Corner Site — Wollstonecraft",
        "p1": [
            ("1A Type",
             "6–7 storey stepped mid-rise · reads lower from upslope"),
            ("1B Character",
             "Garden suburb → apartment renewal · sandstone retaining "
             "walls and mature trees are character-defining"),
            ("1C Precinct",
             "Read as a precinct corner anchor even on a single site — "
             "shapes both Shirley Rd and Romsey St"),
        ],
        "p2": [
            ("2A primary",  "Stepped 5–6–7 storey reading lower upslope"),
            ("2B envelope", "Real envelope < control prism (street-wall step at L5)"),
            ("2C height",   "LEP 21 m · effective height ~17.5 m"),
            ("2D FSR",      "2.0:1 ≈ 3,640 m² · courtyard option (solar evidence)"),
            ("2E depth",    "12–14 m · through-corner replaces deep units"),
            ("2F separation", "12–25 m band: 18 / 12 / 9 m"),
            ("2G street",   "5 m Shirley · 4 m Romsey · sandstone retaining rebuild"),
            ("2H side/rear","6 m + upper-level step · future-neighbour test critical"),
            ("NCC overlay", "Type A · effective height in 12–25 m band"),
        ],
        "p3": [
            "3A  250 mm grid level survey · sandstone outcrop NW retained · 2 fig trees on Romsey verge",
            "3B  long axis N–S · west façade balcony depth + screening",
            "3C  corner lobby at high point · 1:20 accessible ramp + steps",
            "3D  sky terrace at L5 step-back + south ground garden",
            "3E  ~15 % deep soil concentrated downslope (south)",
            "3F  sloping-site overlooking checked in section",
            "3G  continuous accessible route from Romsey footpath",
            "3H  basement entry from Romsey (lower) — minimises ramp",
            "3J  1 sp/unit · 2-level basement · EV ventilation + FRL",
        ],
        "p4": [
            "4A  70 % solar by stepping south side of courtyard down 1 storey",
            "4B  60 % cross-vent · corridors capped 15 m, daylit ends",
            "4C  2.7 m typical · 3.0 m top floor",
            "4D  8.4 m structural grid · 15 / 40 / 35 / 10 mix",
            "4F  short-run corridors · daylit ends",
            "4H  wet stacks aligned · bedrooms off lift & chute",
            "4M  sandstone-toned base ties to retained outcrop",
            "4N  PV + plant screened inside 21 m envelope",
            "4Q  ≥ 20 % adaptable units across orientations",
        ],
        "design": (
            "Site is the evidence: effective height, fire services, "
            "accessible entry and privacy all change because the corner "
            "site falls strongly across the block. Stepped massing keeps "
            "the building under 25 m and protects neighbour amenity."
        ),
        "approval": (
            "Survey + sandstone outcrop plan · stepped-massing 3D · "
            "future-neighbour envelope test (south) · slope-section "
            "privacy diagrams · NCC Type A FRL schedule · fire-engineering "
            "brief · BASIX · landscape & deep-soil · EV ventilation. "
            "PS risk: > 25 m creep + sloping-site privacy."
        ),
    },
    {
        "n": "03",
        "title": "Heritage-Adjacent Infill — Glebe",
        "p1": [
            ("1A Type",
             "3-storey + recessed roof-form terrace-scale infill"),
            ("1B Character",
             "Victorian parapet rhythm · 5.5 m party-wall cadence · "
             "sandstone kerbs — character is the primary control"),
            ("1C Precinct",
             "HCA DCP overrides generic apartment massing under "
             "Housing SEPP s.149 displacement"),
        ],
        "p2": [
            ("2A primary",  "Heritage CMP + DCP rule the form, not the LEP prism"),
            ("2B envelope", "Real envelope much smaller than R1 prism would suggest"),
            ("2C height",   "LEP 11 m · DCP eaves ~6.5 m · 3rd storey set back ≥ 4 m"),
            ("2D FSR",      "0.9:1 ≈ 1,035 m² · ~10 boutique units"),
            ("2E depth",    "11 m glass-to-glass · single rear-court apartment / floor"),
            ("2F separation", "< 12 m band 12/9/6 met internally"),
            ("2G street",   "Inherit 3 m terrace setback exactly · cast-iron palisade"),
            ("2H side/rear","0 m party walls · 6 m rear deep-soil garden"),
            ("NCC overlay", "Type C · effective height < 12 m"),
        ],
        "p3": [
            "3A  heritage CMP · archaeological assessment · party-wall survey #20 + #28",
            "3B  long axis E–W · south-rear unfavourable · living rooms north",
            "3C  restored carriage opening as front entry · CMP palette",
            "3D  rear-garden communal space ~25 % of site · level entry",
            "3E  35 % rear deep soil · tank + plant relocated to roof void",
            "3F  raised sills / offset windows protect #20 first-floor windows",
            "3G  level threshold + accessible ramp behind front fence",
            "3H  no basement (heritage subgrade + small footprint)",
            "3J  1.5 bike sp/unit + pram store · zero on-site car parking",
        ],
        "p4": [
            "4A  ~50 % solar — best-practicable evidence required",
            "4B  dual-aspect every unit → 100 % cross-vent (strong response)",
            "4C  3.0 m ground (Victorian datum) · 2.7 m above",
            "4D  6 × 2B + 4 × 3B family-suited mix",
            "4E  Juliet balconies front · generous rear balconies",
            "4M  reads as one terrace · top floor invisible from footpath",
            "4O  restored rear garden with retained mature tree",
            "4R  existing-fabric upgrade documented in section",
            "4U  BASIX via envelope quality + roof PV · no west glazing",
        ],
        "design": (
            "Heritage character is the primary envelope control. The "
            "building succeeds by reading as one terrace rhythm, with a "
            "recessed attic and rear-garden amenity rather than generic "
            "compliance with R1 prism."
        ),
        "approval": (
            "Heritage Impact Statement · archaeology · party-wall structural "
            "protection · CMP-aligned material schedule · sightline analysis "
            "to listed #20 · best-practicable solar memo · BASIX · NCC Type "
            "C FRL · Council parking-strategy memo."
        ),
    },
    {
        "n": "04",
        "title": "Large Amalgamated Renewal Site — Macquarie Park",
        "p1": [
            ("1A Type",
             "2 slender residential towers above a 4-storey podium with "
             "retail + childcare at ground"),
            ("1B Character",
             "Emerging high-rise renewal precinct · public-domain rules "
             "dominate (street tree, awning, podium height)"),
            ("1C Precinct",
             "Site is large enough to create the public realm — "
             "through-site link, public plaza and green spine"),
        ],
        "p2": [
            ("2A primary",  "2 towers (max plate ~750 m²) + 4-storey podium"),
            ("2B envelope", "Podium 16 m + towers ~78 m · 6 m tower setback off podium edge"),
            ("2C height",   "LEP 80 m · effective height ~75 m → top-band fire matrix"),
            ("2D FSR",      "4.5:1 ≈ 37,800 m² · retail 3,500 + childcare 800 + res ~33,500"),
            ("2E depth",    "Tower plate 18 m max · centre core · dual-aspect units"),
            ("2F separation", "> 25 m band: 24 / 18 / 12 m · tower-to-tower 24 m"),
            ("2G street",   "Podium to street + awning · towers set back 6 m"),
            ("2H side/rear","6 m podium · tower step-back 6 m · future-neighbour critical"),
            ("NCC overlay", "Type A · Class 2 + 6 + 9b · effective height > 50 m, all systems Required"),
        ],
        "p3": [
            "3A  contamination delineation · geotech · wind study · neighbour shadow · traffic",
            "3B  long axis NE–SW · maximise N · podium awning + tower setback for footpath comfort",
            "3C  through-site link · retail spine · plaza outside childcare",
            "3D  L4 podium-roof terrace ~1,500 m² · shaded + sunny zones",
            "3E  ≥ 7 % deep soil concentrated at ground link + tree pockets",
            "3F  24 m tower-to-tower + sky-view test from neighbours",
            "3G  separate residential lobbies / childcare / retail",
            "3H  loading + waste from rear lane · resident basement off side street",
            "3J  3-level basement · 1 sp/unit · end-of-trip · bike near each lobby",
        ],
        "p4": [
            "4A  ≥ 70 % solar by-unit · neighbour shadow is the political risk",
            "4B  centre-core tower · 75 % cross-vent (corner + through)",
            "4C  2.7 m residential · 3.3 m podium retail",
            "4D  10 / 40 / 35 / 15 mix · podium townhouse units",
            "4F  short corridors from centre core · daylit lift lobbies",
            "4H  structural break res ↔ podium · floating floor over plant",
            "4J  acoustic façade to Talavera · winter-garden balconies road side",
            "4M / 4N  modulated tower + warm masonry podium · long-life materials",
            "4O / 4P  engineered podium-roof landscape · soil + drainage cells",
            "4U / 4V / 4W / 4X  all-electric + PV · WSUD · tri-stream waste · BMU",
        ],
        "design": (
            "This is a precinct-scale evidence problem: tower slenderness, "
            "wind, shadow, mixed-use separation, podium landscape and "
            "maintenance access must be coordinated so the public realm, "
            "not the FSR, drives the form."
        ),
        "approval": (
            "Masterplan · contamination + remediation · geotech · wind "
            "tunnel · traffic · neighbour shadow · NCC Type A schedule · "
            "fire-engineering brief · BASIX + Net Zero · podium landscape "
            "engineering · BMU · RAB Act regulated declarations."
        ),
    },
    {
        "n": "05",
        "title": "Noisy Frontage — Strathfield / Parramatta Road",
        "p1": [
            ("1A Type",
             "6-storey shop-top mixed use — retail at ground, apartments above"),
            ("1B Character",
             "Parramatta Rd corridor · auto-oriented · weak pedestrian "
             "environment · character being replaced by renewal"),
            ("1C Precinct",
             "Individual-site with corridor logic — align street wall + "
             "awning continuity to emerging frontage"),
        ],
        "p2": [
            ("2A primary",  "Noise dictates plan + façade — courtyard plan with bedrooms inward"),
            ("2B envelope", "4-storey street wall + 2-storey set back 5 m"),
            ("2C height",   "LEP 21 m · effective height ~18 m → 12–25 m band"),
            ("2D FSR",      "2.5:1 ≈ 5,500 m² · retail 600 + res ~4,900"),
            ("2E depth",    "14 m glass-to-glass · central courtyard plate workable"),
            ("2F separation", "Internal courtyard ≥ 12 m"),
            ("2G street",   "Zero to Parramatta Rd · awning + recessed res entry"),
            ("2H side/rear","6 m rear · future-envelope test on flank"),
            ("NCC overlay", "Type A · Class 2 + 6 · effective height in 12–25 m band"),
        ],
        "p3": [
            "3A  24-hr acoustic survey · NOx air quality · contamination · neighbour windows + POS",
            "3B  front faces north — solar opportunity but noise risk → wintergarden",
            "3C  continuous awning · set-back residential entry · planted glazing buffer",
            "3D  rear courtyard sheltered from road + L4 podium roof terrace",
            "3E  small ground deep-soil + engineered podium-roof landscape (4P)",
            "3F  rear units overlook neighbour yards — angled sightlines + planting",
            "3G  separate residential entry from retail · lobby fully internal",
            "3H  rear lane only — NO Parramatta Rd crossover",
            "3J  0.7 sp/unit · bike room near rear lobby · retail loading off lane",
        ],
        "p4": [
            "4A  N solar opportunity but open-window living constrained",
            "4B  CANNOT count road-facing units cross-vented · realistic ~40 %",
            "4D  bedrooms turned to courtyard or rear · wintergardens road side",
            "4E  wintergardens on Parramatta Rd · open balconies on courtyard",
            "4H  internal acoustic — full-height separation retail ↔ residential L1",
            "4J  acoustic façade to AS/NZS 2107 · filtered mech fresh-air · acoustic refuges",
            "4M  street-side façade engineered for acoustic + thermal · courtyard façade lighter",
            "4S  retail ↔ residential separated by floor + acoustic + entry",
            "4U  mech vent + heat recovery · BASIX via envelope + services strategy",
        ],
        "design": (
            "North-facing solar is not enough. Road-facing windows cannot "
            "be counted for cross-ventilation where noise and air quality "
            "require them to remain closed — the building must perform "
            "with mechanical fresh-air and an acoustic façade."
        ),
        "approval": (
            "24-hr acoustic + air quality · acoustic façade specification · "
            "mechanical ventilation strategy · contamination + remediation · "
            "BASIX · NCC Type A FRL + Part A6 mixed-class · fire-engineering "
            "brief · evidenced cross-vent count · neighbour-impact privacy."
        ),
    },
    {
        "n": "06",
        "title": "Town-Centre Mixed Use — Burwood",
        "p1": [
            ("1A Type",
             "8-storey shop-top — retail + small commercial L1 + apartments above"),
            ("1B Character",
             "Town-centre commercial · shopfront rhythm + awning + "
             "signage hierarchy is the precinct character"),
            ("1C Precinct",
             "Town-centre corner site — reinforce footpath continuity "
             "and contribute to the public realm"),
        ],
        "p2": [
            ("2A primary",  "4-storey street wall + 4-storey upper element set back 4 m"),
            ("2B envelope", "DCP street-wall + Burwood Park sun-access plane shape upper massing"),
            ("2C height",   "LEP 28 m · effective height ~24 m (close to 25 m boundary)"),
            ("2D FSR",      "4.0:1 ≈ 6,600 m² · retail 700 + commercial 600 + res ~5,300"),
            ("2E depth",    "13–14 m double-loaded · corner planning + small light wells"),
            ("2F separation", "12–25 m band: 18 / 12 / 9 m"),
            ("2G street",   "Zero to Burwood Rd · continuous awning at 3.6 m"),
            ("2H side/rear","Side party wall · plan for shared light wells with adjoining redev"),
            ("NCC overlay", "Type A · Class 2 + 5/6 · effective height just under 25 m"),
        ],
        "p3": [
            "3A  survey + Burwood Park solar study · retail tenancy depth · lane logistics",
            "3B  living rooms N + E · upper setback shaped for park afternoon sun",
            "3C  continuous Burwood Rd awning · res entry on lane · refurbished side paving",
            "3D  L5 podium roof terrace ~350 m² · supplemented by Burwood Park 200 m walk",
            "3E  minimal ground deep soil · podium-roof landscape (4P)",
            "3F  internal light-well privacy via offset windows + sill heights",
            "3G  separated residential lobby on lane · retail entries on Burwood Rd",
            "3H  loading + basement off side lane only · NO Burwood Rd crossover",
            "3J  2-level basement · 0.6 sp/unit · ample bike + commercial end-of-trip",
        ],
        "p4": [
            "4A  ≥ 70 % units 2 hr mid-winter sun · NO new park shadow 12:00–14:00 mid-winter",
            "4B  65 % cross-vent · corridors capped 12 m · daylit ends",
            "4C  retail 4.0 / commercial 3.3 / residential 2.7 m (3.0 m at top)",
            "4D  10 / 45 / 35 / 10 mix · 3B units on quieter lane façade",
            "4H  structural break podium ↔ residential · no late-trading uses adjoining",
            "4M  modulated upper façade above masonry podium",
            "4N  PV + plant + communal terrace · plant screened to keep < 25 m",
            "4S / 4T  separate entries · DCP-aligned awning + signage hierarchy",
            "4U / 4V / 4W / 4X  all-electric + PV · greywater podium irrigation · waste swept-path · BMU",
        ],
        "design": (
            "The design must keep the town-centre frontage continuous "
            "while pulling the upper mass away from the park solar plane "
            "and managing the effective-height threshold so Burwood Park "
            "and the town centre both win."
        ),
        "approval": (
            "Burwood Park solar-protection study · town-centre DCP frontage "
            "compliance · NCC Type A + Part A6 schedule · fire-engineering "
            "brief on the 25 m boundary · BASIX · podium landscape "
            "engineering · waste vehicle swept-path · RAB Act declarations."
        ),
    },
]


# --- Drawing helpers ----------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            run = p.add_run()
        else:
            np = tf.add_paragraph()
            np.alignment = align
            run = np.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullet_list(slide, x, y, w, h, items, *, size=10.5,
                    color=INK, font='Calibri', leading=1.18):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_table_rows(slide, x, y, w, rows, *, key_color=GOLD, val_color=INK,
                   size=10.5, row_h=Pt(15)):
    """Render rows as text lines: key (gold) tab value (ink). Uses a textbox with
    two-run paragraphs to avoid PowerPoint's heavyweight table object."""
    h = row_h * (len(rows) + 1)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, (k, v) in enumerate(rows):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.18
        rk = p.add_run()
        rk.text = f"{k}   "
        rk.font.name = "Calibri"
        rk.font.size = Pt(size)
        rk.font.bold = True
        rk.font.color.rgb = key_color
        rv = p.add_run()
        rv.text = v
        rv.font.name = "Calibri"
        rv.font.size = Pt(size)
        rv.font.color.rgb = val_color
    return box


def add_panel(slide, x, y, w, h, title, *, accent=GOLD, num=None):
    """A cream panel with a coloured header strip and section title."""
    # body
    add_rect(slide, x, y, w, h, PAPER, line=RULE)
    # left accent stripe
    add_rect(slide, x, y, Emu(70000), h, accent)
    # title
    add_text(slide, x + Emu(160000), y + Emu(110000),
             w - Emu(200000), Pt(20),
             title, size=10, bold=True, color=NAVY)
    if num:
        add_text(slide, x + w - Emu(700000), y + Emu(110000),
                 Emu(600000), Pt(20),
                 num, size=10, bold=True, color=accent, align=PP_ALIGN.RIGHT)


# --- Slide builder ------------------------------------------------------

# Slide is 16x9 inches = 14,630,400 x 8,229,600 EMU
SW = Inches(16)
SH = Inches(9)

def build_detail_slide(prs, case):
    blank_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(blank_layout)

    # Cream background
    add_rect(slide, 0, 0, SW, SH, CREAM)

    # --- HEADER strip -------------------------------------------------
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.4), Inches(0.13), Inches(8), Inches(0.3),
             f"EXERCISE {case['n']}   |   PERFORMANCE-BASED APARTMENT DESIGN GUIDE",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, Inches(11), Inches(0.13), Inches(4.6), Inches(0.3),
             "INDICATIVE TEACHING EXAMPLE — DETAIL",
             size=10, bold=True, color=HEAD_GOLD, align=PP_ALIGN.RIGHT)

    # --- TITLE area ---------------------------------------------------
    add_text(slide, Inches(0.4), Inches(0.7), Inches(12), Inches(0.55),
             case['title'],
             size=26, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(1.22), Inches(12), Inches(0.32),
             "Detailed evidence breakdown · Parts 1–4 + NCC overlay",
             size=12, color=INK_SOFT)
    # gold rule
    add_rect(slide, Inches(0.4), Inches(1.62), Inches(15.2), Pt(2), GOLD)

    # right-side meta tag mirroring the original page's badge column
    add_rect(slide, Inches(13.0), Inches(0.78), Inches(2.6), Inches(0.7), NAVY)
    add_text(slide, Inches(13.0), Inches(0.85), Inches(2.6), Inches(0.3),
             "DETAILED EVIDENCE", size=9, bold=True,
             color=HEAD_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(13.0), Inches(1.13), Inches(2.6), Inches(0.3),
             f"PAGE 02 / CASE {case['n']}", size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # --- BODY: 4 panels for Parts 1-4 ---------------------------------
    top = Inches(1.85)
    body_h = Inches(5.05)
    gap = Inches(0.18)
    margin_x = Inches(0.4)
    panel_w = (SW - margin_x*2 - gap*3) / 4
    px = margin_x

    # Part 1 panel (Identifying the context)
    add_panel(slide, px, top, panel_w, body_h,
              "PART 1 · IDENTIFYING THE CONTEXT",
              accent=GOLD, num="1A · 1B · 1C")
    inner_x = px + Inches(0.2)
    inner_y = top + Inches(0.55)
    inner_w = panel_w - Inches(0.4)
    for k, v in case['p1']:
        add_text(slide, inner_x, inner_y, inner_w, Pt(13),
                 k, size=10, bold=True, color=GOLD)
        inner_y += Pt(15)
        add_text(slide, inner_x, inner_y, inner_w, Inches(1.6),
                 v, size=10.5, color=INK)
        # estimate vertical advance based on rough wrap
        approx_lines = max(1, (len(v) // 40) + 1)
        inner_y += Pt(13.2 * approx_lines + 6)

    px += panel_w + gap

    # Part 2 panel (Developing the controls)
    add_panel(slide, px, top, panel_w, body_h,
              "PART 2 · DEVELOPING THE CONTROLS",
              accent=GOLD, num="2A → 2H + NCC")
    add_table_rows(slide,
                   px + Inches(0.2), top + Inches(0.55),
                   panel_w - Inches(0.4),
                   case['p2'][:-1],
                   key_color=GOLD, val_color=INK,
                   size=10, row_h=Pt(15.5))
    # NCC overlay row separated as a coloured pill at bottom of panel
    ncc_h = Inches(0.78)
    ncc_y = top + body_h - ncc_h - Inches(0.08)
    add_rect(slide, px + Inches(0.2), ncc_y,
             panel_w - Inches(0.4), ncc_h, NAVY)
    ncc_label, ncc_val = case['p2'][-1]
    add_text(slide, px + Inches(0.32), ncc_y + Inches(0.08),
             panel_w - Inches(0.6), Pt(12),
             ncc_label, size=9, bold=True, color=HEAD_GOLD)
    add_text(slide, px + Inches(0.32), ncc_y + Inches(0.24),
             panel_w - Inches(0.6), ncc_h - Inches(0.30),
             ncc_val, size=10, bold=True, color=WHITE)

    px += panel_w + gap

    # Part 3 panel (Siting)
    add_panel(slide, px, top, panel_w, body_h,
              "PART 3 · SITING THE DEVELOPMENT",
              accent=GOLD, num="3A → 3J")
    add_bullet_list(slide,
                    px + Inches(0.2), top + Inches(0.55),
                    panel_w - Inches(0.4), body_h - Inches(0.7),
                    case['p3'], size=10, leading=1.18)

    px += panel_w + gap

    # Part 4 panel (Designing the building)
    add_panel(slide, px, top, panel_w, body_h,
              "PART 4 · DESIGNING THE BUILDING",
              accent=GOLD, num="4A → 4X")
    add_bullet_list(slide,
                    px + Inches(0.2), top + Inches(0.55),
                    panel_w - Inches(0.4), body_h - Inches(0.7),
                    case['p4'], size=10, leading=1.18)

    # --- Bottom row: teal Design intent + dark red Approval evidence --
    row_top = Inches(7.05)
    row_h = Inches(1.55)
    box_w = (SW - margin_x*2 - Inches(0.18)) / 2

    # Teal box
    add_rect(slide, margin_x, row_top, box_w, row_h, TEAL)
    add_text(slide, margin_x + Inches(0.25), row_top + Inches(0.13),
             box_w - Inches(0.5), Pt(15),
             "DESIGN INTENT ACROSS THE PROJECT",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, margin_x + Inches(0.25), row_top + Inches(0.45),
             box_w - Inches(0.5), row_h - Inches(0.55),
             case['design'], size=12, color=WHITE)

    # Dark red box
    rx = margin_x + box_w + Inches(0.18)
    add_rect(slide, rx, row_top, box_w, row_h, RED)
    add_text(slide, rx + Inches(0.25), row_top + Inches(0.13),
             box_w - Inches(0.5), Pt(15),
             "APPROVAL EVIDENCE PACKAGE",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, rx + Inches(0.25), row_top + Inches(0.45),
             box_w - Inches(0.5), row_h - Inches(0.55),
             case['approval'], size=12, color=WHITE)

    # --- FOOTER -------------------------------------------------------
    add_rect(slide, 0, Inches(8.7), SW, Inches(0.3), NAVY_DARK)
    add_text(slide, Inches(0.4), Inches(8.74), Inches(11), Inches(0.22),
             "Illustrative NSW teaching example only — confirm current "
             "LEP/DCP, Housing SEPP / ADG, NCC 2022, BASIX, DBP and RAB Act "
             "requirements before lodgement.",
             size=8, color=HEAD_GOLD)
    add_text(slide, Inches(11.5), Inches(8.74), Inches(4.1), Inches(0.22),
             "369 Alliance · detailed evidence breakdown",
             size=8, bold=True, color=HEAD_GOLD, align=PP_ALIGN.RIGHT)

    return slide


# --- Slide reordering helper -------------------------------------------

def interleave_slides(prs):
    """Reorder so that newly-appended slides land directly after each
    matching original case slide. Original order: 1..6 (cases). New order
    appended: 7..12 (details). Final order: 1, 7, 2, 8, 3, 9, 4, 10, 5, 11, 6, 12.
    """
    sldIdLst = prs.slides._sldIdLst   # CT_SlideIdList
    children = list(sldIdLst)
    assert len(children) == 12, f"expected 12 slides, got {len(children)}"
    originals = children[:6]
    details = children[6:]
    new_order = []
    for o, d in zip(originals, details):
        new_order.append(o)
        new_order.append(d)
    # Replace in XML element
    for c in children:
        sldIdLst.remove(c)
    for c in new_order:
        sldIdLst.append(c)


# --- Main ---------------------------------------------------------------

def main():
    src = r'C:\Users\raque\Downloads\ADG_6_Worked_Site_Exercise_Realistic_Drawings.pptx'
    out = r'C:\2026\369 Alliance\Website\docs\ADG_6_Worked_Site_Exercise_With_Detail_Pages.pptx'
    prs = Presentation(src)
    print(f'Loaded {src}')
    print(f'  Slides: {len(prs.slides)}  Size: {prs.slide_width}x{prs.slide_height} EMU')

    for c in CASES:
        build_detail_slide(prs, c)
        print(f'  + Detail page for case {c["n"]} · {c["title"]}')

    interleave_slides(prs)
    prs.save(out)
    print(f'\nSaved {out}')
    print(f'  Final slide count: {len(prs.slides)}')


if __name__ == '__main__':
    main()
