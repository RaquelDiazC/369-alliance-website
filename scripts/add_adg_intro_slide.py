"""
Insert a single introductory slide at position 1 of
ADG_6_Worked_Site_Exercise_With_Detail_Pages.pptx, holding the 6-case
overview table (Address + 'why it's a different exercise').
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette (sampled from existing slides)
NAVY      = RGBColor(0x14, 0x25, 0x36)
NAVY_DARK = RGBColor(0x20, 0x32, 0x44)
CREAM     = RGBColor(0xfe, 0xfc, 0xf8)
PAPER     = RGBColor(0xff, 0xff, 0xff)
RULE      = RGBColor(0xd8, 0xcd, 0xb6)
GOLD      = RGBColor(0xc9, 0xa9, 0x6d)
TEAL      = RGBColor(0x3e, 0x79, 0x7d)
RED       = RGBColor(0x9c, 0x5a, 0x56)
INK       = RGBColor(0x14, 0x25, 0x36)
INK_SOFT  = RGBColor(0x4a, 0x55, 0x60)
HEAD_GOLD = RGBColor(0xe8, 0xd8, 0xae)
WHITE     = RGBColor(0xff, 0xff, 0xff)
ROW_ALT   = RGBColor(0xf6, 0xef, 0xde)


CASES = [
    ("01", "Inner West Flat Infill",      "Marrickville",
     "14–18 Premier St, Marrickville NSW 2204",
     "Flat infill · low-rise · NCC Type C baseline — DtS pathway viable",
     "Type C · < 12 m"),
    ("02", "Sloping Corner Site",         "Wollstonecraft",
     "3–7 Shirley Rd cnr Romsey St, Wollstonecraft NSW 2065",
     "Sloping corner · mid-rise · slope-section privacy test",
     "Type A · ~17.5 m"),
    ("03", "Heritage-Adjacent Infill",    "Glebe",
     "22–26 Mitchell St, Glebe NSW 2037",
     "Heritage Conservation Area — character displaces ADG metrics",
     "Type C · < 12 m"),
    ("04", "Large Amalgamated Renewal",   "Macquarie Park",
     "100–120 Talavera Rd, Macquarie Park NSW 2113",
     "Amalgamated podium + 2 towers · effective height ~75 m",
     "Type A · ~75 m"),
    ("05", "Noisy Frontage Shop-Top",     "Strathfield",
     "245–251 Parramatta Rd, Strathfield NSW 2135",
     "Noisy / polluted frontage — cross-ventilation count constrained",
     "Type A · ~18 m"),
    ("06", "Town-Centre Mixed Use",       "Burwood",
     "30–34 Burwood Rd, Burwood NSW 2134",
     "Town-centre shop-top · sitting on the 25 m fire-matrix boundary",
     "Type A · ~24 m"),
]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            run = p.add_run()
        else:
            np = tf.add_paragraph()
            np.alignment = align
            run = np.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def build_intro_slide(prs):
    SW = prs.slide_width
    SH = prs.slide_height
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background
    add_rect(slide, 0, 0, SW, SH, CREAM)

    # Top header strip
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.4), Inches(0.13), Inches(8), Inches(0.3),
             "OVERVIEW   |   PERFORMANCE-BASED APARTMENT DESIGN GUIDE",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, Inches(11), Inches(0.13), Inches(4.6), Inches(0.3),
             "INDICATIVE TEACHING EXAMPLES — INDEX",
             size=10, bold=True, color=HEAD_GOLD, align=PP_ALIGN.RIGHT)

    # Title
    add_text(slide, Inches(0.4), Inches(0.7), Inches(12), Inches(0.55),
             "Six worked NSW site exercises",
             size=26, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(1.22), Inches(12), Inches(0.32),
             "Each case is deliberately different — the deck stress-tests "
             "different parts of the guide and different NCC overlays.",
             size=12, color=INK_SOFT)
    add_rect(slide, Inches(0.4), Inches(1.62), Inches(15.2), Pt(2), GOLD)

    # Right-side meta tag
    add_rect(slide, Inches(13.0), Inches(0.78), Inches(2.6), Inches(0.7), NAVY)
    add_text(slide, Inches(13.0), Inches(0.85), Inches(2.6), Inches(0.3),
             "TEACHING DECK", size=9, bold=True,
             color=HEAD_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(13.0), Inches(1.13), Inches(2.6), Inches(0.3),
             "6 CASES · 12 PAGES", size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # ----- TABLE -----
    margin_x = Inches(0.4)
    table_top = Inches(1.95)
    table_w = SW - margin_x*2

    # Column widths (sum = table_w)
    col_n_w     = Inches(0.55)
    col_loc_w   = Inches(2.4)
    col_addr_w  = Inches(4.3)
    col_why_w   = Inches(5.95)
    col_ncc_w   = table_w - col_n_w - col_loc_w - col_addr_w - col_why_w

    header_h = Inches(0.42)
    row_h    = Inches(0.78)

    # Header row
    add_rect(slide, margin_x, table_top, table_w, header_h, NAVY)
    hx = margin_x
    for w, label in [
        (col_n_w,    "#"),
        (col_loc_w,  "LOCATION"),
        (col_addr_w, "ADDRESS (ILLUSTRATIVE)"),
        (col_why_w,  "WHY IT'S A DIFFERENT EXERCISE"),
        (col_ncc_w,  "NCC / HEIGHT"),
    ]:
        add_text(slide, hx + Inches(0.18), table_top + Inches(0.10),
                 w - Inches(0.2), header_h - Inches(0.1),
                 label, size=10, bold=True, color=HEAD_GOLD,
                 align=PP_ALIGN.LEFT if label != "#" else PP_ALIGN.CENTER)
        hx += w

    # Body rows
    y = table_top + header_h
    for i, (n, title, locality, addr, why, ncc) in enumerate(CASES):
        # Row background
        bg = ROW_ALT if i % 2 == 0 else PAPER
        add_rect(slide, margin_x, y, table_w, row_h, bg, line=RULE)
        # Left gold accent for the case number cell
        add_rect(slide, margin_x, y, col_n_w, row_h, GOLD)
        # # column
        add_text(slide, margin_x, y + Inches(0.20),
                 col_n_w, Inches(0.4),
                 n, size=18, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)

        cx = margin_x + col_n_w
        # Location column (title + locality)
        add_text(slide, cx + Inches(0.18), y + Inches(0.10),
                 col_loc_w - Inches(0.2), Inches(0.30),
                 title, size=10.5, bold=True, color=NAVY)
        add_text(slide, cx + Inches(0.18), y + Inches(0.40),
                 col_loc_w - Inches(0.2), Inches(0.30),
                 locality, size=10.5, color=GOLD, bold=True)
        cx += col_loc_w

        # Address
        add_text(slide, cx + Inches(0.18), y + Inches(0.22),
                 col_addr_w - Inches(0.2), Inches(0.5),
                 addr, size=11, color=INK)
        cx += col_addr_w

        # Why (the focus column)
        add_text(slide, cx + Inches(0.18), y + Inches(0.22),
                 col_why_w - Inches(0.2), Inches(0.5),
                 why, size=11, bold=True, color=NAVY)
        cx += col_why_w

        # NCC pill
        pill_x = cx + Inches(0.18)
        pill_y = y + Inches(0.20)
        pill_w = col_ncc_w - Inches(0.36)
        pill_h = Inches(0.40)
        add_rect(slide, pill_x, pill_y, pill_w, pill_h, NAVY)
        add_text(slide, pill_x, pill_y + Inches(0.08),
                 pill_w, pill_h - Inches(0.1),
                 ncc, size=10.5, bold=True, color=HEAD_GOLD,
                 align=PP_ALIGN.CENTER)

        y += row_h

    # ----- Bottom row: teal + dark red strap matching the design -----
    row_top = Inches(7.05)
    row_h2  = Inches(1.55)
    box_w   = (SW - margin_x*2 - Inches(0.18)) / 2

    add_rect(slide, margin_x, row_top, box_w, row_h2, TEAL)
    add_text(slide, margin_x + Inches(0.25), row_top + Inches(0.13),
             box_w - Inches(0.5), Pt(15),
             "FIVE-PASS WORKFLOW APPLIED TO EVERY CASE",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, margin_x + Inches(0.25), row_top + Inches(0.45),
             box_w - Inches(0.5), row_h2 - Inches(0.55),
             "1 site intelligence  →  2 massing strategy  →  "
             "3 apartment performance  →  4 regulated design  →  "
             "5 occupation + operation. Every page in this deck inherits "
             "the same workflow logic.",
             size=12, color=WHITE)

    rx = margin_x + box_w + Inches(0.18)
    add_rect(slide, rx, row_top, box_w, row_h2, RED)
    add_text(slide, rx + Inches(0.25), row_top + Inches(0.13),
             box_w - Inches(0.5), Pt(15),
             "DECK STRUCTURE",
             size=10, bold=True, color=HEAD_GOLD)
    add_text(slide, rx + Inches(0.25), row_top + Inches(0.45),
             box_w - Inches(0.5), row_h2 - Inches(0.55),
             "Each case spans two pages: a realistic 3D drawing "
             "with site evidence rail, then a detailed Parts 1–4 + "
             "NCC overlay breakdown. Total 12 case pages + this index.",
             size=12, color=WHITE)

    # Footer
    add_rect(slide, 0, Inches(8.7), SW, Inches(0.3), NAVY_DARK)
    add_text(slide, Inches(0.4), Inches(8.74), Inches(11), Inches(0.22),
             "Illustrative NSW teaching examples only — confirm current "
             "LEP/DCP, Housing SEPP / ADG, NCC 2022, BASIX, DBP and RAB "
             "Act requirements before lodgement.",
             size=8, color=HEAD_GOLD)
    add_text(slide, Inches(11.5), Inches(8.74), Inches(4.1), Inches(0.22),
             "369 Alliance · teaching deck index",
             size=8, bold=True, color=HEAD_GOLD, align=PP_ALIGN.RIGHT)

    return slide


def move_to_front(prs):
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    # Newly added intro is the last one — move it to position 0
    intro = children[-1]
    sldIdLst.remove(intro)
    sldIdLst.insert(0, intro)


def main():
    src = r'C:\2026\369 Alliance\Website\docs\ADG_6_Worked_Site_Exercise_With_Detail_Pages.pptx'
    out = r'C:\2026\369 Alliance\Website\docs\ADG_6_Worked_Site_Exercise_With_Detail_Pages.pptx'
    prs = Presentation(src)
    print(f'Loaded {src} ({len(prs.slides)} slides)')
    build_intro_slide(prs)
    move_to_front(prs)
    prs.save(out)
    print(f'Saved {out} ({len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
